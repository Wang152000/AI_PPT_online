import streamlit as st
import io
import chardet
from docx import Document
import requests
from bs4 import BeautifulSoup
import re
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import ChatOpenAI
from langchain.prompts import PromptTemplate
from langchain.chains import LLMChain
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import tempfile
import os

# 设置页面
st.set_page_config(
    page_title="智能PPT生成器",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# 自定义CSS样式
st.markdown("""
<style>
.main {
    padding: 0rem 1rem;
}
.stTextInput>div>div>input {
    padding: 0.5rem;
}
.content-box {
    border: 1px solid #ddd;
    padding: 20px;
    border-radius: 5px;
    background-color: #f8f9fa;
    margin: 10px 0;
}
.step-title {
    color: #1976D2;
    font-size: 1.5em;
    font-weight: bold;
    margin-bottom: 20px;
    padding-bottom: 10px;
    border-bottom: 2px solid #1976D2;
}
.editor-container {
    border: 1px solid #ddd;
    border-radius: 5px;
    padding: 10px;
    margin: 10px 0;
    background-color: white;
}
.editor-toolbar {
    padding: 10px;
    background-color: #f8f9fa;
    border-bottom: 1px solid #ddd;
    margin-bottom: 10px;
}
.word-count {
    color: #666;
    font-size: 0.9em;
    text-align: right;
    padding: 5px;
}
.button-container {
    display: flex;
    justify-content: space-between;
    margin: 10px 0;
}
.article-display {
    font-size: 16px;
    line-height: 1.8;
    white-space: pre-wrap;
    font-family: "Microsoft YaHei", sans-serif;
    padding: 20px;
    background-color: white;
    border-radius: 5px;
    border: 1px solid #ddd;
    margin: 20px 0;
}
.article-display p {
    margin-bottom: 1em;
}
.stTextArea textarea {
    font-size: 16px !important;
    line-height: 1.8 !important;
    font-family: "Microsoft YaHei", sans-serif !important;
}
.comparison-box {
    border: 1px solid #ddd;
    border-radius: 5px;
    padding: 10px;
    margin: 10px 0;
    background-color: #f8f9fa;
}
.content-title {
    font-size: 1.2em;
    font-weight: bold;
    margin-bottom: 10px;
    padding-bottom: 5px;
    border-bottom: 2px solid #1976D2;
    color: #1976D2;
}
</style>
""", unsafe_allow_html=True)

def detect_encoding(file_content):
    """检测文件编码"""
    result = chardet.detect(file_content)
    return result['encoding']

def extract_text_from_txt(file):
    """从txt文件中提取文本"""
    try:
        content = file.read()
        encoding = detect_encoding(content)
        text = content.decode(encoding)
        return text
    except Exception as e:
        return f"错误：无法读取TXT文件。原因：{str(e)}"

def extract_text_from_docx(file):
    """从docx文件中提取文本"""
    try:
        # 检查文件大小
        file.seek(0, 2)
        file_size = file.tell()
        file.seek(0)
        
        if file_size == 0:
            return "错误：文件为空。请确保上传了有效的Word文档。"
            
        file_content = file.read()
        file_in_memory = io.BytesIO(file_content)
        
        try:
            doc = Document(file_in_memory)
        except Exception as doc_error:
            if "There is no item named 'NULL' in the archive" in str(doc_error):
                return "错误：文件格式不正确。请确保：\n1. 文件是真正的.docx格式（不是重命名的.doc文件）\n2. 文件未被损坏\n3. 文件不是空白文档"
            else:
                return f"错误：无法读取DOCX文件。原因：{str(doc_error)}"
        
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
        
        text = '\n\n'.join(full_text)
        
        if not text:
            return "错误：文档内容为空。请确保文档包含文本内容。"
            
        return text
        
    except Exception as e:
        error_msg = str(e)
        if "Permission denied" in error_msg:
            return "错误：无法访问文件。请确保文件未被其他程序占用。"
        elif "not a zip file" in error_msg.lower():
            return "错误：文件格式不正确。请确保上传的是正确的.docx格式文件。"
        else:
            return f"错误：无法读取DOCX文件。原因：{error_msg}"

def extract_article_from_url(url):
    """从URL中提取文章内容"""
    try:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        
        if response.encoding == 'ISO-8859-1':
            response.encoding = response.apparent_encoding
        
        soup = BeautifulSoup(response.text, 'html.parser')
        
        # 移除不需要的标签
        for script in soup(["script", "style", "meta", "link", "header", "footer", "nav"]):
            script.decompose()
        
        # 查找可能的文章内容容器
        article_containers = soup.find_all(['article', 'div'], class_=re.compile(r'article|content|post|text|body'))
        
        if not article_containers:
            # 如果找不到特定容器，尝试查找最长的文本块
            paragraphs = soup.find_all('p')
            if paragraphs:
                # 保留段落的原始格式
                text_blocks = []
                for p in paragraphs:
                    if len(p.get_text().strip()) > 50:
                        # 处理换行
                        text = p.get_text().strip()
                        # 保留br标签的换行
                        for br in p.find_all('br'):
                            br.replace_with('\n')
                        text_blocks.append(text)
                text = '\n\n'.join(text_blocks)
            else:
                return "错误：无法从该网页提取有效的文章内容。"
        else:
            # 使用找到的最长的文章容器
            main_container = max(article_containers, key=lambda x: len(x.get_text()))
            
            # 保留原始格式
            text_blocks = []
            
            # 处理段落和换行
            for element in main_container.descendants:
                if element.name == 'p' and element.get_text().strip():
                    # 处理段落内的换行
                    text = element.get_text().strip()
                    # 保留br标签的换行
                    for br in element.find_all('br'):
                        br.replace_with('\n')
                    text_blocks.append(text)
                elif element.name == 'br':
                    text_blocks.append('\n')
                elif element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                    # 保留标题格式
                    text_blocks.append(f"\n\n{element.get_text().strip()}\n")
                elif element.name == 'li':
                    # 保留列表项格式
                    text_blocks.append(f"• {element.get_text().strip()}\n")
            
            text = '\n'.join(text_blocks)
        
        # 清理文本但保留有意义的换行
        text = re.sub(r'\n{3,}', '\n\n', text)  # 将3个以上的换行减少为2个
        text = re.sub(r' {2,}', ' ', text)  # 删除多余的空格
        
        if len(text) < 100:
            return "错误：提取的文本内容过短，可能不是有效的文章。"
        
        return text
        
    except requests.RequestException as e:
        return f"错误：无法访问该URL。原因：{str(e)}"
    except Exception as e:
        return f"错误：提取文章内容失败。原因：{str(e)}"

def recursive_split_text(text, num_chunks):
    """使用递归字符分割文本，基于指定的块数进行分割"""
    try:
        # 计算每个块的大致大小
        total_length = len(text)
        chunk_size = total_length // num_chunks

        # 确保chunk_size不会太小
        chunk_size = max(chunk_size, 100)

        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=0,
            length_function=len,
            is_separator_regex=False,
            separators=["\n\n", "\n", "。", "！", "？", ".", "!", "?", " ", ""]
        )

        chunks = text_splitter.split_text(text)

        # 调整块数以匹配目标数量
        if len(chunks) > num_chunks:
            # 如果块数过多，合并相邻的块
            while len(chunks) > num_chunks:
                # 找到最短的相邻块对
                min_length = float('inf')
                merge_index = 0
                for i in range(len(chunks) - 1):
                    combined_length = len(chunks[i]) + len(chunks[i + 1])
                    if combined_length < min_length:
                        min_length = combined_length
                        merge_index = i

                # 合并这对块
                chunks[merge_index] = chunks[merge_index] + chunks[merge_index + 1]
                chunks.pop(merge_index + 1)

        elif len(chunks) < num_chunks:
            # 如果块数不足，分割最长的块
            while len(chunks) < num_chunks:
                # 找到最长的块
                max_length = 0
                split_index = 0
                for i, chunk in enumerate(chunks):
                    if len(chunk) > max_length:
                        max_length = len(chunk)
                        split_index = i

                # 分割最长的块
                chunk_to_split = chunks[split_index]
                split_point = len(chunk_to_split) // 2

                # 在句子边界处分割
                separators = ["。", "！", "？", ".", "!", "?", "\n"]
                best_split_point = split_point
                min_distance = float('inf')

                # 在分割点附近寻找最近的分隔符
                for i, char in enumerate(chunk_to_split):
                    if char in separators:
                        distance = abs(i - split_point)
                        if distance < min_distance:
                            min_distance = distance
                            best_split_point = i + 1  # 包含分隔符

                # 如果没有找到合适的分隔符，就使用原始分割点
                if best_split_point == split_point:
                    # 确保不会在单词中间分割
                    while (best_split_point < len(chunk_to_split) and
                           chunk_to_split[best_split_point].isalnum()):
                        best_split_point += 1
                    while (best_split_point > 0 and
                           chunk_to_split[best_split_point-1].isalnum()):
                        best_split_point -= 1

                # 执行分割
                chunks[split_index] = chunk_to_split[:best_split_point]
                chunks.insert(split_index + 1, chunk_to_split[best_split_point:])

        return chunks

    except Exception as e:
        st.error(f"递归分割失败：{str(e)}")
        return None

def extract_content(text_block, api_key, base_url):
    """使用大模型提炼文本内容并生成标题"""
    try:
        llm = ChatOpenAI(
            openai_api_key=api_key,
            openai_api_base=base_url,
            temperature=0.7,
            model_name="gpt-3.5-turbo"
        )

        prompt_template = """##目标
提取并总结输入内容的关键信息，形成层次分明的要点说明，同时生成一个简短的标题（不超过20个字）。

##要求：
（1）内容完整性：
- 保持原文的主要内容和关键信息，在原文基础上适度精简
- 保留重要的数据、案例和专业术语
- 确保每个要点都有充分的解释和必要的上下文

（2）层级结构：
- 识别并保持原文的层级关系
- 使用缩进表示不同层级（每个子层级缩进2个空格）
- 保持原文的逻辑组织结构
- 对并列关系、递进关系、因果关系等进行清晰的层级划分

（3）格式规范：
- 使用数字编号标识主要层级（1. 2. 3.）
- 使用字母编号标识次级层级（a. b. c.）
- 使用符号标识更深层级（- 或 •）
- 每个层级的标题使用3-8个字的短语
- 在标题后详细展开该层级的具体内容
- 使用分号分隔复杂内容中的多个方面

（4）表达方式：
- 保持专业性和准确性
- 使用清晰、简洁的语言
- 避免过度概括和模糊表达
- 保留原文的重要表述方式和专业用语

##特别说明：
即使原文已经包含分点内容，也必须重新组织和提炼，确保内容更加精炼和结构化。

##输入
{text_block}

##输出格式
标题：[简短的标题]

内容：
1. [一级标题]：
  a. [二级要点]：[详细说明]
    - [三级要点]：[具体内容]
  b. [二级要点]：[详细说明]
2. [一级标题]：
  a. [二级要点]：[详细说明]
    - [三级要点]：[具体内容]
……

注意：
1. 严格遵守缩进规则，确保层级关系清晰
2. 保持原文的重要细节和专业表述
3. 适度精简但不过度概括
4. 确保每个层级都有充分的说明和解释
"""

        prompt = PromptTemplate(
            template=prompt_template,
            input_variables=["text_block"]
        )

        chain = LLMChain(llm=llm, prompt=prompt)
        result = chain.invoke({
            "text_block": text_block
        })

        # 解析结果，分离标题和内容
        output_text = result['text']
        title = ""
        content = ""
        
        # 分离标题和内容
        lines = output_text.split('\n')
        for i, line in enumerate(lines):
            if line.startswith('标题：'):
                title = line.replace('标题：', '').strip()
            elif line.startswith('内容：'):
                content = '\n'.join(lines[i+1:]).strip()
                break

        return content, title, False  # 返回提炼内容、标题和一个标志表示这不是分点内容

    except Exception as e:
        st.error(f"内容提炼失败：{str(e)}")
        return None, None, False

def generate_main_title(extracted_contents):
    """基于全文内容生成总标题"""
    try:
        # 收集所有文本内容
        all_content = ""
        for item in extracted_contents:
            all_content += item['title'] + "\n" + item['content'] + "\n\n"

        # 使用LLM生成总标题
        llm = ChatOpenAI(
            openai_api_key=st.session_state['api_key'],
            openai_api_base=st.session_state['base_url'],
            temperature=0.7,
            model_name="gpt-3.5-turbo"
        )

        prompt_template = """请基于以下文章内容，生成一个简短的总标题（不超过20个字）。标题应该：
1. 准确概括文章的核心主题
2. 使用简洁有力的语言
3. 避免过于笼统的表述
4. 突出文章的独特性和价值

文章内容：
{text}

请直接输出标题，不要添加任何其他内容。"""

        prompt = PromptTemplate(
            template=prompt_template,
            input_variables=["text"]
        )

        chain = LLMChain(llm=llm, prompt=prompt)
        result = chain.invoke({
            "text": all_content
        })

        return result['text'].strip()
    except Exception as e:
        st.error(f"生成总标题失败：{str(e)}")
        return "内容提炼报告"

def create_ppt(extracted_contents):
    """创建PPT文件"""
    prs = Presentation()
    
    # 设置幻灯片尺寸为16:9
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # 创建封面
    cover_slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # 生成总标题
    main_title = generate_main_title(extracted_contents)
    
    # 设置主标题
    title = cover_slide.shapes.title
    title.text = main_title
    title.text_frame.paragraphs[0].font.size = Pt(60)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 118, 210)  # 使用蓝色
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 设置副标题
    subtitle = cover_slide.placeholders[1]
    subtitle.text = "内容提炼报告"
    subtitle.text_frame.paragraphs[0].font.size = Pt(40)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 118, 210)  # 使用蓝色
    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 为每个提炼内容创建幻灯片
    for item in extracted_contents:
        # 创建新的幻灯片（使用空白布局）
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 使用完全空白的布局
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = item['title']
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.color.rgb = RGBColor(31, 118, 210)  # 使用蓝色
        title_para.font.bold = True
        
        # 添加内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(6.5))
        content_frame = content_box.text_frame
        
        # 解析并添加内容
        lines = item['content'].split('\n')
        for line in lines:
            if line.strip():
                p = content_frame.add_paragraph()
                p.text = line.strip()
                
                # 设置字体格式
                p.font.size = Pt(18)  # 统一使用18号字
                p.line_spacing = 1.5  # 设置1.5倍行距
                
                if line.strip().startswith(('1.', '2.', '3.', '4.', '5.')):  # 一级标题
                    p.font.bold = True
                    p.font.size = Pt(28)  # 一级标题保持28号字
                elif line.strip().startswith(('a.', 'b.', 'c.', 'd.')):  # 二级要点
                    p.font.bold = True
                    p.level = 1
                elif line.strip().startswith(('-', '•')):  # 三级要点
                    p.font.bold = True
                    p.level = 2
                else:  # 普通内容
                    p.level = 3
    
    # 保存PPT
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
        prs.save(tmp.name)
        return tmp.name

def main():
    """主函数"""
    # 初始化session state
    if 'step' not in st.session_state:
        st.session_state['step'] = 1
    if 'extracted_text' not in st.session_state:
        st.session_state['extracted_text'] = None
    if 'edited_text' not in st.session_state:
        st.session_state['edited_text'] = None
    if 'is_editing' not in st.session_state:
        st.session_state['is_editing'] = False
    if 'chunks' not in st.session_state:
        st.session_state['chunks'] = None
    if 'edited_chunks' not in st.session_state:
        st.session_state['edited_chunks'] = []
    if 'extracted_contents' not in st.session_state:
        st.session_state['extracted_contents'] = []
    if 'api_key_confirmed' not in st.session_state:
        st.session_state['api_key_confirmed'] = False
    if 'block_operations' not in st.session_state:
        st.session_state['block_operations'] = {'insert_index': None}

    # 设置页面标题和样式
    st.title("智能PPT生成器")
    
    # 显示当前步骤
    st.markdown(f"### 当前步骤：{st.session_state['step']}/3")
    
    # 根据步骤显示不同的页面
    if st.session_state['step'] == 1:
        show_step1()
    elif st.session_state['step'] == 2:
        show_step2()
    elif st.session_state['step'] == 3:
        show_step3()

def show_step1():
    """显示第一步：文件上传和内容提取"""
    st.markdown('<div class="step-box">', unsafe_allow_html=True)
    st.markdown('<div class="step-title">步骤1：上传文件</div>', unsafe_allow_html=True)
    st.markdown("""
    请上传Word文档（.docx格式）或文本文件（.txt格式），或者输入文章URL。系统将自动提取文档内容。

    **提示：**
    - 支持的文件格式：.docx, .txt
    - 支持直接输入URL地址
    - 文件大小限制：10MB
    """, unsafe_allow_html=True)
    st.markdown('</div>', unsafe_allow_html=True)

    # 选择输入方式
    input_method = st.radio(
        "选择输入方式",
        ["上传文件", "输入URL"],
        horizontal=True
    )

    if input_method == "上传文件":
        # 文档上传部分
        uploaded_file = st.file_uploader(
            "上传文档",
            type=['txt', 'docx'],
            help="支持的文件格式：TXT、DOCX"
        )

        if uploaded_file:
            if st.button("提取文章"):
                with st.spinner('正在提取文章内容...'):
                    if uploaded_file.name.endswith('.txt'):
                        text = extract_text_from_txt(uploaded_file)
                    elif uploaded_file.name.endswith('.docx'):
                        text = extract_text_from_docx(uploaded_file)
                    else:
                        text = "错误：不支持的文件格式"

                    st.session_state['extracted_text'] = text
                    st.session_state['edited_text'] = text

    else:  # 输入URL
        url = st.text_input("输入文章URL", help="请输入包含文章的网页地址")
        if url and st.button("提取文章"):
            with st.spinner('正在从URL提取文章内容...'):
                text = extract_article_from_url(url)
                st.session_state['extracted_text'] = text
                st.session_state['edited_text'] = text

    # 显示提取的文章内容
    if st.session_state['extracted_text']:
        st.write("### 文章内容")
        
        # 如果正在编辑
        if st.session_state['is_editing']:
            # 创建文本编辑器
            edited_text = st.text_area(
                "编辑文章内容",
                value=st.session_state['edited_text'],
                height=400
            )
            
            # 字数统计
            word_count = len(edited_text)
            st.markdown(f"<div class='word-count'>字数：{word_count}</div>", unsafe_allow_html=True)
            
            # 保存和取消按钮
            col1, col2 = st.columns(2)
            with col1:
                if st.button("保存修改"):
                    st.session_state['edited_text'] = edited_text
                    st.session_state['is_editing'] = False
                    st.rerun()
            with col2:
                if st.button("取消修改"):
                    st.session_state['is_editing'] = False
                    st.session_state['edited_text'] = st.session_state['extracted_text']
                    st.rerun()
        else:
            # 显示文章内容
            st.markdown(f"<div class='article-display'>{st.session_state['edited_text']}</div>", unsafe_allow_html=True)
            
            # 编辑按钮
            col1, col2 = st.columns(2)
            with col1:
                if st.button("编辑文章"):
                    st.session_state['is_editing'] = True
                    st.rerun()
            with col2:
                if st.button("确认内容并进入下一步"):
                    st.session_state['step'] = 2
                    st.rerun()

def show_step2():
    """显示第二步：内容分割"""
    st.markdown('<div class="step-box">', unsafe_allow_html=True)
    st.markdown('<div class="step-title">步骤2：内容分割</div>', unsafe_allow_html=True)
    st.markdown("""
    系统将自动将文档内容分割成多个文本块。您可以：
    - 调整分割块数
    - 编辑每个文本块的内容
    - 删除不需要的文本块
    - 在任意位置添加新的文本块
    - 确认分割结果后进入下一步

    **提示：**
    - 拖动滑块调整分割块数
    - 点击编辑按钮修改文本块内容
    - 点击删除按钮移除不需要的文本块
    - 使用"添加新块"按钮在任意位置插入新的文本块
    """, unsafe_allow_html=True)
    st.markdown('</div>', unsafe_allow_html=True)

    # 分割参数设置
    num_chunks = st.slider(
        "分割块数",
        min_value=2,
        max_value=20,
        value=5,
        step=1,
        help="将文章分割成几个部分"
    )

    if st.button("应用分割", key="split_button"):
        with st.spinner('正在进行文本分割...'):
            chunks = recursive_split_text(
                st.session_state['edited_text'],
                num_chunks
            )
            if chunks:
                st.session_state['chunks'] = chunks
                st.session_state['edited_chunks'] = chunks.copy()

    # 显示分割结果
    if st.session_state['chunks']:
        st.write("### 分割预览")
        
        # 使用容器增加宽度
        with st.container():
            # 在开头添加"新增文章块"按钮
            if st.button("在开头添加新块 ⬆", key="insert_start"):
                st.session_state['block_operations']['insert_index'] = 0
                st.rerun()

            # 处理插入操作
            if st.session_state['block_operations']['insert_index'] is not None:
                idx = st.session_state['block_operations']['insert_index']
                if 0 <= idx <= len(st.session_state['edited_chunks']):
                    st.session_state['edited_chunks'].insert(idx, "在这里输入新的内容...")
                    st.session_state['block_operations']['insert_index'] = None
                    st.rerun()

            # 显示所有块
            for i, chunk in enumerate(st.session_state['edited_chunks']):
                st.markdown(f"#### 第 {i+1} 部分")
                
                # 计算所需的高度：每行25像素，额外加50像素作为缓冲
                num_lines = len(chunk.split('\n'))
                height = max(num_lines * 25 + 50, 200)  # 最小高度为200像素
                
                # 创建可编辑的文本区域
                edited_text = st.text_area(
                    "",
                    value=chunk,
                    height=height,
                    key=f"chunk_{i}"
                )
                st.session_state['edited_chunks'][i] = edited_text

                # 操作按钮行
                col1, col2 = st.columns([1, 9])
                with col1:
                    # 删除按钮
                    if len(st.session_state['edited_chunks']) > 1:  # 保持至少一个块
                        if st.button("🗑️", key=f"delete_{i}", help="删除此块"):
                            st.session_state['edited_chunks'].pop(i)
                            st.rerun()

                # 在每个块之后添加"新增文章块"按钮
                if st.button(f"在此处添加新块 ⬇", key=f"insert_{i}"):
                    st.session_state['block_operations']['insert_index'] = i + 1
                    st.rerun()

                st.markdown("---")

            # 在末尾添加"新增文章块"按钮
            if st.button("在末尾添加新块 ⬇", key="insert_end"):
                st.session_state['block_operations']['insert_index'] = len(st.session_state['edited_chunks'])
                st.rerun()

        # 操作按钮
        col1, col2 = st.columns(2)
        with col1:
            if st.button("返回上一步"):
                st.session_state['step'] = 1
                st.session_state['chunks'] = None
                st.session_state['edited_chunks'] = []
                st.rerun()
        
        with col2:
            if st.button("确认分割并进入下一步"):
                st.session_state['step'] = 3
                st.rerun()

def show_step3():
    """显示第三步：内容提炼和PPT生成"""
    st.markdown('<div class="step-box">', unsafe_allow_html=True)
    st.markdown('<div class="step-title">步骤3：内容提炼</div>', unsafe_allow_html=True)
    st.markdown("""
    系统将对每个文本块进行内容提炼，并生成结构化的PPT内容。

    **提示：**
    - 每个文本块都会生成对应的PPT页面
    - 您可以预览生成的PPT效果
    - 确认无误后可以导出PPT文件
    """, unsafe_allow_html=True)
    st.markdown('</div>', unsafe_allow_html=True)

    # 检查是否有文本块需要处理
    if not st.session_state.get('edited_chunks'):
        st.warning("没有找到需要处理的文本块，请返回上一步添加内容。")
        if st.button("返回上一步"):
            st.session_state['step'] = 2
            st.rerun()
        return

    # API设置部分
    with st.expander("API设置", expanded=True):
        base_url = st.text_input(
            "API基础URL",
            value=st.session_state.get('base_url', "https://api.gpt.ge/v1/"),
            help="请输入API基础URL",
            key="base_url_input"
        )
        api_key = st.text_input(
            "API密钥",
            type="password",
            value=st.session_state.get('api_key', ''),
            help="请输入您的API密钥",
            key="api_key_input"
        )
        
        # 添加确认和重置按钮
        if api_key:
            col1, col2 = st.columns([1, 1])
            with col1:
                if not st.session_state.get('api_key_confirmed', False):
                    if st.button("确认API密钥", key="confirm_api_key"):
                        st.session_state['api_key'] = api_key
                        st.session_state['base_url'] = base_url
                        st.session_state['api_key_confirmed'] = True
                        st.rerun()
            with col2:
                if st.session_state.get('api_key_confirmed', False):
                    if st.button("重置API密钥", key="reset_api_key"):
                        st.session_state['api_key'] = ''
                        st.session_state['api_key_confirmed'] = False
                        st.rerun()
            
            if not st.session_state.get('api_key_confirmed', False):
                st.info("请点击确认按钮以验证API密钥")
            else:
                st.success("API密钥已确认，可以开始内容提炼")

    # 内容提炼部分
    if st.session_state.get('api_key') and st.session_state.get('api_key_confirmed', False):
        if not st.session_state.get('extracted_contents'):
            if st.button("开始内容提炼"):
                progress_bar = st.progress(0)
                status_text = st.empty()

                # 存储提炼结果
                extracted_contents = []
                total_chunks = len(st.session_state['edited_chunks'])

                for i, chunk in enumerate(st.session_state['edited_chunks']):
                    status_text.text(f"正在处理第 {i+1}/{total_chunks} 个文本块...")
                    progress_bar.progress((i + 1) / total_chunks)

                    try:
                        content, title, is_points = extract_content(
                            chunk,
                            st.session_state['api_key'],
                            st.session_state['base_url']
                        )
                        
                        if content and title:
                            extracted_contents.append({
                                'title': title,
                                'content': content,
                                'original': chunk
                            })
                    except Exception as e:
                        st.error(f"处理文本块时发生错误：{str(e)}")
                        break

                if extracted_contents:
                    st.session_state['extracted_contents'] = extracted_contents
                    status_text.empty()
                    progress_bar.empty()
                    st.rerun()
                else:
                    st.error("内容提炼失败，请检查API密钥是否正确或重试")

        # 显示提炼结果
        if st.session_state.get('extracted_contents'):
            st.write("### 内容提炼预览")
            
            for i, item in enumerate(st.session_state['extracted_contents']):
                st.markdown(f"#### 第 {i+1} 部分：{item['title']}")
                
                # 使用列布局创建左右对照
                col1, col2 = st.columns(2)
                
                with col1:
                    st.markdown('<div class="comparison-box">', unsafe_allow_html=True)
                    st.markdown('<div class="content-title">原文内容</div>', unsafe_allow_html=True)
                    st.markdown(f"<div class='article-display' style='height: 400px; overflow-y: auto;'>{item['original']}</div>", unsafe_allow_html=True)
                    st.markdown('</div>', unsafe_allow_html=True)
                
                with col2:
                    st.markdown('<div class="comparison-box">', unsafe_allow_html=True)
                    st.markdown('<div class="content-title">提炼结果</div>', unsafe_allow_html=True)
                    st.markdown(f"<div class='article-display' style='height: 400px; overflow-y: auto;'>{item['content']}</div>", unsafe_allow_html=True)
                    st.markdown('</div>', unsafe_allow_html=True)
                
                st.markdown("---")

            # 添加导出PPT按钮
            if st.button("导出为PPT"):
                with st.spinner("正在生成PPT..."):
                    try:
                        # 创建PPT文件
                        ppt_path = create_ppt(st.session_state['extracted_contents'])
                        
                        # 读取文件内容
                        with open(ppt_path, 'rb') as file:
                            ppt_data = file.read()
                        
                        # 提供下载按钮
                        st.download_button(
                            label="下载PPT文件",
                            data=ppt_data,
                            file_name="content_summary.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                        
                        # 删除临时文件
                        os.unlink(ppt_path)
                        
                    except Exception as e:
                        st.error(f"生成PPT时发生错误：{str(e)}")

    elif not st.session_state.get('api_key'):
        st.error("请先在API设置中输入API密钥")

    # 操作按钮
    col1, col2 = st.columns(2)
    with col1:
        if st.button("返回上一步"):
            st.session_state['step'] = 2
            st.session_state['extracted_contents'] = []
            st.rerun()
    
    with col2:
        if st.button("重新开始"):
            # 重置所有状态
            st.session_state['step'] = 1
            st.session_state['extracted_text'] = None
            st.session_state['edited_text'] = None
            st.session_state['is_editing'] = False
            st.session_state['chunks'] = None
            st.session_state['edited_chunks'] = []
            st.session_state['extracted_contents'] = []
            st.session_state['api_key_confirmed'] = False
            st.rerun()

if __name__ == "__main__":
    main() 