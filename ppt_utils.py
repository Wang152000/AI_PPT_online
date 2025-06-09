from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import streamlit as st
import io
import re

def get_indent_level(line):
    """获取行的缩进级别"""
    # 计算行首空格数量
    leading_spaces = len(line) - len(line.lstrip())
    # 每2个空格算作一个缩进级别
    return leading_spaces // 2

def format_line(line):
    """格式化行内容，处理特殊字符和格式"""
    # 移除行首的空格
    line = line.lstrip()
    # 处理不同类型的标识符
    if re.match(r'^\d+\.', line):  # 数字编号
        return line, True, 0
    elif re.match(r'^[a-z]\.', line):  # 字母编号
        return line, True, 1
    elif line.startswith('- ') or line.startswith('• '):  # 符号标识
        return line, True, 2
    return line, False, 0

def create_slide(prs, title, content):
    """创建一个新的PPT幻灯片，支持层级缩进"""
    # 使用标题和内容布局
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # 设置内容
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()  # 清除默认文本
    text_frame.word_wrap = True
    
    # 处理每一行内容
    lines = content.split('\n')
    for line in lines:
        if not line.strip():
            continue
            
        # 获取缩进级别和格式化行
        indent_level = get_indent_level(line)
        formatted_line, is_bullet, bullet_level = format_line(line)
        
        # 创建新段落
        p = text_frame.add_paragraph()
        p.text = formatted_line
        
        # 设置字体
        p.font.size = Pt(18)  # 基础字体大小
        if is_bullet and bullet_level == 0:  # 一级标题
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(31, 73, 125)  # 深蓝色
        elif is_bullet and bullet_level == 1:  # 二级标题
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(68, 84, 106)  # 灰蓝色
        
        # 设置缩进
        p.level = indent_level
        if indent_level > 0:
            # 左侧缩进，每级缩进0.5英寸
            p.left_indent = Inches(0.5 * indent_level)
        
        # 设置行间距
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        
        # 设置对齐方式
        p.alignment = PP_ALIGN.LEFT
    
    return slide

def preview_ppt_in_streamlit(extracted_contents):
    """在Streamlit中预览PPT内容"""
    st.write("### PPT预览")
    
    # 添加自定义CSS样式
    st.markdown("""
        <style>
        .ppt-preview {
            border: 1px solid #ddd;
            padding: 20px;
            border-radius: 5px;
            background-color: white;
            margin: 10px 0;
        }
        .ppt-title {
            font-size: 24px;
            font-weight: bold;
            color: #1f497d;
            margin-bottom: 15px;
            padding-bottom: 10px;
            border-bottom: 2px solid #eee;
        }
        .ppt-content {
            font-size: 16px;
            line-height: 1.6;
        }
        .indent-1 { margin-left: 20px; }
        .indent-2 { margin-left: 40px; }
        .indent-3 { margin-left: 60px; }
        </style>
    """, unsafe_allow_html=True)
    
    for i, content in enumerate(extracted_contents):
        with st.expander(f"第 {i+1} 页：{content['title']}", expanded=True):
            st.markdown('<div class="ppt-preview">', unsafe_allow_html=True)
            st.markdown(f'<div class="ppt-title">{content["title"]}</div>', unsafe_allow_html=True)
            
            # 处理内容的层级显示
            content_lines = content['content'].split('\n')
            formatted_lines = []
            for line in content_lines:
                indent_level = get_indent_level(line)
                if indent_level > 0:
                    formatted_lines.append(f'<div class="indent-{indent_level}">{line.lstrip()}</div>')
                else:
                    formatted_lines.append(line)
            
            st.markdown(
                '<div class="ppt-content">' + 
                '\n'.join(formatted_lines) + 
                '</div>',
                unsafe_allow_html=True
            )
            st.markdown('</div>', unsafe_allow_html=True)

def export_ppt(extracted_contents):
    """导出PPT文件"""
    # 创建演示文稿
    prs = Presentation()
    
    # 设置幻灯片大小为16:9
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # 添加封面
    cover_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = cover_slide.shapes.title
    title.text = "内容提炼报告"
    subtitle = cover_slide.placeholders[1]
    subtitle.text = f"共 {len(extracted_contents)} 页"
    
    # 为每个内容块创建幻灯片
    for content in extracted_contents:
        create_slide(prs, content['title'], content['content'])
    
    # 保存到内存中
    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    
    return ppt_buffer.getvalue() 